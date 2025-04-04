from langchain_core.tools import Tool
from langchain_google_community import GoogleSearchAPIWrapper
from functools import partial
import os
from pptx import Presentation
from my_mcp.mcp_server import mcp

@mcp.tool(name="google_search",
          description="Search Google for recent results")
def google_search_web_tool(num_results=5):
    """
    Searches Google for recent results and returns a callable function that performs the search.

    Parameters:
        num_results (int): The number of search results to retrieve. Must be a positive integer.

    Returns:
        callable: A function that when executed returns search results from Google.

    Raises:
        TypeError: If 'num_results' is not an integer.
        ValueError: If 'num_results' is not a positive integer.
    """
    if not isinstance(num_results, int):
        raise TypeError("num_results must be an integer")
    if num_results <= 0:
        raise ValueError("num_results must be a positive integer")

    search = GoogleSearchAPIWrapper()
    func = partial(search.results, num_results=num_results)
    return func

@mcp.tool(name="create_ppt",
          description="Generates a PowerPoint presentation with provided title, content and references.")
def generate_presentation_tool(title: str, content: str, references: list):
    """
    Generates a PowerPoint presentation with a title slide and a references slide.

    The presentation includes:
      - A title slide displaying the provided 'title' and 'content'.
      - A references slide listing all items in the 'references' list. If the list is empty,
        the slide will indicate that no references were found.

    Parameters:
        title (str): The title for the presentation. Must be a non-empty string.
        content (str): The content for the title slide.
        references (list): A list of reference strings. Each item should be a string.

    Returns:
        str: A message indicating the file path where the presentation was saved.

    Raises:
        TypeError: If 'title' or 'content' is not a string or if 'references' is not a list.
        ValueError: If 'title' is empty.
        Exception: Propagates any exception raised during the creation or saving of the presentation.
    """
    if not isinstance(title, str):
        raise TypeError("title must be a string")
    if not title.strip():
        raise ValueError("title cannot be empty")
    if not isinstance(content, str):
        raise TypeError("content must be a string")
    if not isinstance(references, list):
        raise TypeError("references must be a list")

    try:
        prs = Presentation()

        # Create title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = content

        # Create references slide
        if len(prs.slide_layouts) > 1:
            ref_slide_layout = prs.slide_layouts[1]
        else:
            ref_slide_layout = title_slide_layout

        ref_slide = prs.slides.add_slide(ref_slide_layout)
        ref_slide.shapes.title.text = "References"
        if references:
            ref_text = "\n".join(references)
        else:
            ref_text = "No references found."
        if len(ref_slide.placeholders) > 1:
            ref_slide.placeholders[1].text = ref_text

        # Save the presentation to the current working directory
        ppt_path = os.path.join(os.getcwd(), "presentation.pptx")
        prs.save(ppt_path)
        return f"Presentation saved to: {ppt_path}"
    except Exception as e:
        raise Exception(f"Failed to generate presentation: {e}")

@mcp.tool(name="add_references",
          description="Processes and formats a list of references.")
def process_references_tool(references: list):
    """
    Processes and formats a list of references into a single string output.

    Parameters:
        references (list): A list of reference strings.

    Returns:
        str: A formatted string containing all processed references. If the list is empty,
             a message indicating that no references are available is returned.

    Raises:
        TypeError: If 'references' is not a list.
    """
    if not isinstance(references, list):
        raise TypeError("references must be a list")

    if not references:
        return "No references to process."

    ref_output = "Processed References:\n" + "\n".join(ref for ref in references)
    return ref_output
